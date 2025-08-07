from fastapi import FastAPI, UploadFile, HTTPException, Request, File
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
from pypdf import PdfReader
from openai import OpenAI  # Nueva importación
import faiss
import numpy as np
from dotenv import load_dotenv
import os
import json
from bs4 import BeautifulSoup
from pptx import Presentation
import httpx
import logging
from pydantic import BaseModel
from typing import List, Dict
from urllib.parse import urljoin
from io import BytesIO

# Intentar imports opcionales para extracción avanzada y OCR
try:
    import trafilatura  # type: ignore
    HAS_TRAFILATURA = True
except Exception:
    HAS_TRAFILATURA = False

try:
    from PIL import Image  # type: ignore
    import pytesseract  # type: ignore
    HAS_OCR = True
except Exception:
    HAS_OCR = False

app = FastAPI()

# Configuración de logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuración CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))  # Nueva inicialización

# --- Persistencia con Faiss y JSON ---
FAISS_INDEX_PATH = "vectorstore.faiss"
CHUNKS_PATH = "chunks.json"

# Cargar índice y chunks si existen
if os.path.exists(FAISS_INDEX_PATH):
    index = faiss.read_index(FAISS_INDEX_PATH)
else:
    index = faiss.IndexFlatL2(1536)  # Dimensión de text-embedding-3-small

if os.path.exists(CHUNKS_PATH):
    with open(CHUNKS_PATH, 'r', encoding='utf-8') as f:
        stored_chunks = json.load(f)
else:
    stored_chunks = []
# --- Fin de la configuración de persistencia ---


def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

# --- Extracción robusta de texto de páginas HTML ---
async def extract_text_from_page(html: str, page_url: str, http_client: httpx.AsyncClient) -> str:
    parts: List[str] = []

    # 1) Extraer con trafilatura si está disponible (contenido principal)
    if HAS_TRAFILATURA:
        try:
            extracted = trafilatura.extract(html, url=page_url, include_links=True)
            if extracted:
                parts.append(extracted)
        except Exception:
            pass

    soup = BeautifulSoup(html, 'html.parser')

    # 2) Título y meta
    try:
        if soup.title and soup.title.string:
            parts.append(soup.title.string.strip())
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        if meta_desc and meta_desc.get('content'):
            parts.append(meta_desc['content'].strip())
        og_title = soup.find('meta', property='og:title')
        if og_title and og_title.get('content'):
            parts.append(og_title['content'].strip())
        og_desc = soup.find('meta', property='og:description')
        if og_desc and og_desc.get('content'):
            parts.append(og_desc['content'].strip())
    except Exception:
        pass

    # 3) Texto visible de contenedores
    try:
        visible_text = "\n".join(s for s in soup.stripped_strings if s)
        if visible_text:
            parts.append(visible_text)
    except Exception:
        pass

    # 4) Atributos accesibles alt/title/aria-label/value
    try:
        for el in soup.find_all(True):
            for attr in ('alt', 'title', 'aria-label', 'value'):
                val = el.get(attr)
                if isinstance(val, str) and val.strip():
                    parts.append(val.strip())
    except Exception:
        pass

    # 5) Texto de enlaces, botones y labels
    try:
        for a in soup.find_all('a'):
            txt = a.get_text(strip=True)
            if txt:
                parts.append(txt)
        for b in soup.find_all(['button', 'label']):
            txt = b.get_text(strip=True)
            if txt:
                parts.append(txt)
    except Exception:
        pass

    # 6) OCR opcional en imágenes (limitado)
    if HAS_OCR:
        try:
            img_tags = soup.find_all('img')[:5]
            for img in img_tags:
                src = img.get('src') or ''
                if not src or src.lower().endswith('.svg'):
                    continue
                img_url = urljoin(page_url, src)
                try:
                    resp = await http_client.get(img_url, timeout=10.0)
                    if resp.status_code == 200 and resp.content:
                        try:
                            image = Image.open(BytesIO(resp.content))
                            text_ocr = pytesseract.image_to_string(image)
                            if text_ocr and text_ocr.strip():
                                parts.append(text_ocr.strip())
                        except Exception:
                            continue
                except Exception:
                    continue
        except Exception:
            pass

    # Desduplicar líneas
    combined = "\n".join(parts)
    lines = [ln.strip() for ln in combined.splitlines() if ln and ln.strip()]
    seen = set()
    unique_lines: List[str] = []
    for ln in lines:
        if ln not in seen:
            seen.add(ln)
            unique_lines.append(ln)

    return "\n".join(unique_lines)

def process_and_store_text(text: str):
    """Chunks text, creates embeddings, and stores them."""
    if not text or not text.strip():
        return 0
    
    chunks = chunk_text(text)
    
    if not chunks:
        return 0

    embeddings = []
    for chunk in chunks:
        emb = client.embeddings.create(
            input=chunk,
            model="text-embedding-3-small"
        ).data[0].embedding
        embeddings.append(emb)
    
    if embeddings:
        index.add(np.array(embeddings))
        stored_chunks.extend(chunks)

        # Guardar cambios en disco
        faiss.write_index(index, FAISS_INDEX_PATH)
        with open(CHUNKS_PATH, 'w', encoding='utf-8') as f:
            json.dump(stored_chunks, f, ensure_ascii=False, indent=2)
    
    return len(chunks)

@app.post("/upload_pdf")
async def upload_file(file: UploadFile = File(...)):
    try:
        # Verificar tipo de archivo
        if not file.filename.endswith('.pdf'):
            raise HTTPException(status_code=400, detail="Solo se aceptan archivos PDF")
        
        reader = PdfReader(file.file)
        full_text = ""
        for page in reader.pages:
            try:
                page_text = page.extract_text()
                if page_text:
                    full_text += page_text + "\n"
            except:
                continue
        
        if not full_text.strip():
            raise HTTPException(
            status_code=400,
            detail="No pude leer el texto de este PDF. Por favor verifica que: \n1. El PDF contenga texto seleccionable (no sea una imagen escaneada)\n2. No esté protegido con contraseña\n3. Tenga al menos un párrafo de texto"
            )

        num_chunks = process_and_store_text(full_text)
        
        if num_chunks == 0:
            raise HTTPException(status_code=500, detail="No se pudieron procesar fragmentos del documento.")

        return {
    "status": "success",
    "message": f"PDF procesado. {num_chunks} fragmentos almacenados.",
    "full_text": full_text  # ✅ importante
}
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar PDF: {str(e)}")

@app.post("/upload_pptx")
async def upload_pptx(file: UploadFile):
    try:
        if not file.filename.endswith(('.pptx')):
            raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
        
        full_text = ""
        presentation = Presentation(file.file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"

        if not full_text.strip():
            raise HTTPException(status_code=400, detail="No se encontró texto en el archivo PPTX.")

        num_chunks = process_and_store_text(full_text)

        if num_chunks == 0:
            raise HTTPException(status_code=500, detail="No se pudieron procesar fragmentos del documento.")

        return {
    "status": "success",
    "message": f"PPTX procesado. {num_chunks} fragmentos almacenados.",
    "full_text": full_text
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar PPTX: {str(e)}")

class ExcelRequest(BaseModel):
    data: List[Dict]

@app.post("/upload-excel")
async def upload_excel(request: ExcelRequest):
    try:
        if not request.data:
            return {
                "success": False,
                "message": "No se recibieron datos del archivo Excel."
            }
        
        # Crear texto estructurado para cada fila
        full_text = ""
        processed_rows = 0
        
        for row_data in request.data:
            # Obtener los valores de las primeras 3 columnas/keys
            # El frontend usa XLSX que puede generar keys como __EMPTY, __EMPTY_1, etc.
            # o usar los headers de la primera fila
            
            keys = list(row_data.keys())
            if len(keys) < 3:
                continue  # Saltar filas con menos de 3 columnas
            
            # Obtener pregunta, respuesta y tema (primeras 3 columnas)
            pregunta_key = keys[0]
            respuesta_key = keys[1] 
            tema_key = keys[2]
            
            pregunta = str(row_data.get(pregunta_key, "")).strip()
            respuesta = str(row_data.get(respuesta_key, "")).strip() 
            tema = str(row_data.get(tema_key, "General")).strip()
            
            # Limpiar valores que pueden venir como "nan" o valores vacíos
            if pregunta.lower() in ['nan', 'null', ''] or respuesta.lower() in ['nan', 'null', '']:
                continue
                
            if not tema or tema.lower() in ['nan', 'null', '']:
                tema = "General"
            
            # Solo procesar si hay pregunta y respuesta válidas
            if pregunta and respuesta:
                # Formato estructurado que ayuda al modelo a entender el contexto
                chunk_text = f"""TEMA: {tema}
PREGUNTA: {pregunta}
RESPUESTA: {respuesta}

---"""
                full_text += chunk_text + "\n"
                processed_rows += 1
        
        if not full_text.strip():
            return {
                "success": False,
                "message": "No se encontraron preguntas y respuestas válidas en el archivo Excel."
            }
        
        # Procesar y almacenar el texto
        num_chunks = process_and_store_text(full_text)
        
        if num_chunks == 0:
            return {
                "success": False,
                "message": "No se pudieron procesar fragmentos del documento Excel."
            }
        
        return {
            "success": True,
            "message": f"✅ Excel procesado correctamente. {processed_rows} preguntas-respuestas almacenadas, {num_chunks} fragmentos creados.",
            "processed_rows": processed_rows,
            "total_chunks": num_chunks
        }
        
    except Exception as e:
        logger.error(f"Error al procesar Excel desde frontend: {e}", exc_info=True)
        return {
            "success": False,
            "message": f"Error al procesar el archivo Excel: {str(e)}"
        }

class UrlRequest(BaseModel):
    empresaId: str = None  # Opcional para compatibilidad
    links: list[str] = []  # Acepta el formato del frontend
    urls: list[str] = []   # Mantiene compatibilidad con formato anterior

@app.post("/scrape_url")
async def scrape_url(request: UrlRequest):
    try:
        # Determinar qué lista de URLs usar (compatibilidad con ambos formatos)
        url_list = request.links if request.links else request.urls
        logger.info(f"URLs recibidas: {url_list}")
        
        # Filtrar links vacíos
        valid_links = [link.strip() for link in url_list if link.strip()]
        logger.info(f"URLs válidas después de filtrar: {valid_links}")
        
        if not valid_links:
            logger.warning("No se proporcionaron links válidos")
            return {
                "success": False,
                "message": "No se proporcionaron links válidos."
            }
        
        total_chunks = 0
        processed_urls = []
        failed_urls = []
        
        for url in valid_links:
            try:
                # Agregar https:// si no tiene protocolo
                if not url.startswith(('http://', 'https://')):
                    url = 'https://' + url
                
                logger.info(f"Procesando URL: {url}")
                
                async with httpx.AsyncClient() as http_client:
                    response = await http_client.get(url, follow_redirects=True, timeout=20.0)
                    response.raise_for_status()

                    # Extracción robusta (con OCR si disponible)
                    full_text = await extract_text_from_page(response.text, url, http_client)

                logger.info(f"Texto extraído de {url}: {len(full_text)} caracteres")
                
                if not full_text.strip():
                    logger.warning(f"No se encontró texto procesable en {url}")
                    failed_urls.append({"url": url, "error": "No se encontró texto procesable"})
                    continue

                num_chunks = process_and_store_text(full_text)
                logger.info(f"Creados {num_chunks} chunks para {url}")
                total_chunks += num_chunks
                processed_urls.append({"url": url, "chunks": num_chunks})
                
            except httpx.HTTPStatusError as e:
                logger.error(f"Error HTTP en {url}: {e.response.status_code}")
                failed_urls.append({"url": url, "error": f"Error HTTP {e.response.status_code}"})
            except httpx.RequestError as e:
                logger.error(f"Error de red en {url}: {str(e)}")
                failed_urls.append({"url": url, "error": f"Error de red: {str(e)}"})
            except Exception as e:
                logger.error(f"Error inesperado en {url}: {str(e)}", exc_info=True)
                failed_urls.append({"url": url, "error": f"Error inesperado: {str(e)}"})

        logger.info(f"Resumen del scraping: {len(processed_urls)} URLs exitosas, {len(failed_urls)} URLs fallidas, {total_chunks} chunks totales")
        
        if not processed_urls:
            return {
                "success": False,
                "message": f"No se pudo procesar ningún link. Se intentaron {len(valid_links)} URLs pero todas fallaron. Revisa los detalles en failed_urls.",
                "failed_urls": failed_urls,
                "total_attempted": len(valid_links)
            }

        return {
            "success": True,
            "message": f"✅ {len(processed_urls)} links procesados exitosamente. {total_chunks} fragmentos de información almacenados. {len(failed_urls)} URLs fallaron.",
            "processed_urls": processed_urls,
            "failed_urls": failed_urls,
            "total_chunks": total_chunks,
            "total_attempted": len(valid_links)
        }

    except HTTPException:
        raise
    except Exception as e:
        return {
            "success": False,
            "message": f"Error al procesar los links: {str(e)}"
        }

class QuestionRequest(BaseModel):
    question: str

@app.post("/ask")
async def ask_question(request: Request):
    try:
        # Compatibilidad con diferentes formatos de solicitud
        try:
            data = await request.json()
        except:
            body = await request.body()
            data = json.loads(body)
        
        question = data.get("question")
        if not question:
            logger.info("Solicitud a /ask sin pregunta.")
            return {
                "answer": "Parece que no se ha enviado ninguna pregunta. Por favor, inténtalo de nuevo.",
                "context": []
            }
        
        if index.ntotal == 0:
            logger.info("Solicitud a /ask con índice vacío.")
            return {
                "answer": "Actualmente no tengo documentos con los que trabajar. Por favor, pide a un administrador que suba la información necesaria (PDF, PPTX, etc.) para que pueda responder a tus preguntas.",
                "context": []
            }
            
        # Generar embedding
        q_embedding = client.embeddings.create(  # Llamada actualizada
            input=question,
            model="text-embedding-3-small"
        ).data[0].embedding

        # Búsqueda de chunks similares
        _, I = index.search(np.array([q_embedding]), k=3)
        similar_chunks = [stored_chunks[i] for i in I[0] if i < len(stored_chunks)]
        context = "\n---\n".join(similar_chunks)

        # Generar respuesta
        system_prompt = """Eres un asistente virtual de atención al cliente. Tu misión es ayudar a los usuarios de forma amable y profesional.
1.  **Analiza el contexto**: Basa tu respuesta únicamente en el fragmento de texto proporcionado en el contexto. No utilices conocimiento externo ni inventes información.
2.  **Tono**: Mantén un tono corporativo pero cercano y amigable. Usa un lenguaje claro y evita la jerga técnica si es posible.
3.  **Si encuentras la respuesta**: Proporciónala de manera completa y fácil de entender.
4.  **Si NO encuentras la respuesta**: Indica amablemente que la información no se encuentra en los documentos disponibles. Ofrece alternativas, como contactar a un agente humano o revisar la pregunta.
5.  **Cierre proactivo**: Finaliza siempre tus respuestas de forma positiva, preguntando si hay algo más en lo que puedas ayudar. Por ejemplo: '¿Hay algo más en lo que pueda asistirte hoy?'"""

        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Contexto:\n{context}\n\nPregunta: {question}"}
            ],
            temperature=0.3
        )

        return {
            "answer": response.choices[0].message.content,
            "context": similar_chunks
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error inesperado en /ask: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
